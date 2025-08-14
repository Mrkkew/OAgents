from typing import Optional, Callable, Dict, Any, List
from smolagents import Tool
from smolagents.models import MessageRole, Model
from .mdconvert import MarkdownConverter
from xml.dom import minidom
from openpyxl import load_workbook
from openpyxl.styles import PatternFill
import json
from Bio import PDB
from pptx import Presentation
import os

MAX_ROWS = 500
TEXT_LIMIT_DEFAULT = 100000

class TextInspectorTool(Tool):
    name = "inspect_file_as_text"
    description = """
You cannot load files yourself: instead call this tool to read a file as markdown text and ask questions about it.
This tool handles the following file extensions: [".html", ".pdb", ".xlsx", ".xls", ".pdf", ".docx", ".ppt", ".pptx"], and all other types of text files. IT DOES NOT HANDLE IMAGES."""

    inputs = {
        "file_path": {
            "description": "The path to the file you want to read as text. Must be a '.something' file, like '.pdf'. If it is an image, use the visualizer tool instead! If it is an audio, use the audio tool instead! DO NOT use this tool for an HTML webpage: use the web_search tool instead!",
            "type": "string",
        },
        "question": {
            "description": "[Optional]: Your question, as a natural language sentence. Provide as much context as possible. Do not pass this parameter if you just want to directly return the content of the file.",
            "type": "string",
            "nullable": True,
        },
    }
    output_type = "string"
    SUPPORTED_EXTS = {".html", ".pdb", ".xlsx", ".xls", ".pdf", ".docx", ".ppt", ".pptx"}
    UNSUPPORTED_EXTS = {".png", ".jpg", ".zip"}
    md_converter = MarkdownConverter()

    def __init__(self, model: Model, text_limit: int = TEXT_LIMIT_DEFAULT):
        super().__init__()
        self.model = model
        self.text_limit = text_limit

    def jsonld_to_markdown(self, data):
        markdown = ""
        if isinstance(data, dict):
            for key, value in data.items():
                markdown += f"**{key}**: {self.jsonld_to_markdown(value)}\n"
        elif isinstance(data, list):
            for item in data:
                markdown += f"- {self.jsonld_to_markdown(item)}\n"
        else:
            markdown += str(data)
        return markdown
    
    
    def _process_file(self, file_path):
        ext = os.path.splitext(file_path)[1].lower()
        handlers: Dict[str, Callable[[str], str]] = {
            ".xml": self._parse_xml,
            ".csv": self._parse_csv,
            ".pdb": lambda fp: self._parse_pdb_file(fp),
            ".ppt": self._parse_ppt_file,
            ".pptx": self._parse_ppt_file,
            ".xls": self._parse_excel_file,
            ".xlsx": self._parse_excel_file,
            ".jsonld": self._parse_jsonld,
        }
        if ext in handlers:
            return handlers[ext](file_path)
        elif any(ext in unsupported for unsupported in self.UNSUPPORTED_EXTS):
            raise RuntimeError(f"Cannot use inspect_file_as_text tool with {ext}: use appropriate tool instead!")
        else:
            result = self.md_converter.convert(file_path)
            return result.text_content
    
    def _parse_pdb_file(self, file_path):
        parser = PDB.PDBParser(QUIET=True)
        structure = parser.get_structure("protein", file_path)
        
        atoms = list(structure.get_atoms())
        if len(atoms) < 2:
            return "Error: PDB file contains fewer than two atoms."
        
        atom1, atom2 = atoms[0], atoms[1]
        distance = atom1 - atom2 
        
        return f"First atom: {atom1.get_name()} ({atom1.coord})\n" \
            f"Second atom: {atom2.get_name()} ({atom2.coord})\n" \
            f"Distance: {distance:.3f} Angstroms ({distance * 100:.0f} pm)"

    def _parse_excel_file(self, file_path, max_rows=MAX_ROWS):
        try:
            workbook = load_workbook(filename=file_path, read_only=True)
            all_sheets_text = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                result = []
                
                row_count = 0
                for row in sheet.iter_rows():
                    if row_count >= max_rows:
                        break
                        
                    row_data = []
                    for cell in row:
                        cell_value = cell.value if cell.value is not None else ""        
                        cell_color = "FFFFFF" 
                        try:
                            fill = cell.fill
                            if hasattr(fill, "fgColor") and hasattr(fill.fgColor, "rgb"):
                                rgb = fill.fgColor.rgb
                                if rgb and isinstance(rgb, str) and len(rgb) == 8:
                                    cell_color = rgb[2:]
                        except:
                            pass
                            
                        row_data.append({
                            "value": str(cell_value),
                            "color": cell_color
                        })
                    result.append(row_data)
                    row_count += 1
                    
                sheet_text = []
                num_rows = len(result)
                num_cols = len(result[0]) if result else 0
                sheet_text.append(f"Table '{sheet_name}' contains {num_rows} rows {num_cols} column:")
                for row in result:
                    row_text = ""
                    for cell in row:
                        value = cell["value"] if cell["value"] != "" else "None"
                        color = cell["color"]
                        if color == "FFFFFF" or color == "000000":
                            row_text += f"{value} "
                        else:
                            row_text += f"{value}({color}) "
                    sheet_text.append(row_text)
                
                all_sheets_text.append("\n".join(sheet_text))
            
            return "\n\n".join(all_sheets_text)
        except Exception as e:
            raise RuntimeError(f"Error parsing Excel file: {str(e)}")


    def _parse_ppt_file(self, file_path):
        content = ""
        try:
            ppt = Presentation(file_path)
            for slide_number, slide in enumerate(ppt.slides, start=1):
                content += f"=== Slide {slide_number} ===\n"
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text.strip())
                content += "\n".join(slide_texts) + "\n"
            return content.strip()
        except Exception as e:
            raise RuntimeError(f"Error parsing PPT file: {e}")
        
    def _parse_xml(self, file_path):
        try:
            dom = minidom.parse(file_path)
            texts = [
                node.firstChild.nodeValue
                for node in dom.getElementsByTagName("*")
                if node.firstChild and node.firstChild.nodeType == node.TEXT_NODE
            ]
            return " ".join(texts).strip()
        except Exception as e:
            raise RuntimeError(f"Error parsing XML file: {str(e)}")
        
    def _parse_csv(self, file_path):
        try:
            with open(file_path, "r") as fr:
                contents = fr.readlines()
            return "".join(contents)
        except Exception as e:
            raise RuntimeError(f"Error parsing CSV file: {str(e)}")
        
    def _parse_jsonld(self, file_path):
        try:
            with open(file_path, "r") as f:
                data = json.load(f)
            result_text = self.jsonld_to_markdown(data)
            return result_text
        except Exception as e:
            raise RuntimeError(f"Error parsing JSON-LD file: {str(e)}")
    
    def forward_initial_exam_mode(self, file_path, question):
        try:
            content = self._process_file(file_path)

            if not question:
                return content

            messages = [
                {
                    "role": MessageRole.SYSTEM,
                    "content": [
                        {
                            "type": "text",
                            "text": "Here is a file:\n### "
                            + str(file_path)
                            + "\n"
                            + content[: self.text_limit],
                        }
                    ],
                },
                {
                    "role": MessageRole.USER,
                    "content": [
                        {
                            "type": "text",
                            "text": "Now please write a short, 5 sentence caption for this document, that could help someone asking this question: "
                            + question
                            + "\nDon't answer the question yourself! Just provide useful notes on the document",
                        }
                    ],
                },
            ]
            return self.model(messages).content
        except Exception as e:
            return str(e)

    def forward(self, file_path, question: Optional[str] = None) -> str:
        try:
            content = self._process_file(file_path)

            if not question:
                return content

            messages = [
                {
                    "role": MessageRole.SYSTEM,
                    "content": [
                        {
                            "type": "text",
                            "text": "You will have to write a short caption for this file, then answer this question:"
                            + question,
                        }
                    ],
                },
                {
                    "role": MessageRole.USER,
                    "content": [
                        {
                            "type": "text",
                            "text": "Here is the complete file:\n### "
                            + str(file_path)
                            + "\n"
                            + content[: self.text_limit],
                        }
                    ],
                },
                {
                    "role": MessageRole.USER,
                    "content": [
                        {
                            "type": "text",
                            "text": "Now answer the question below. Use these three headings: '1. Short answer', '2. Extremely detailed answer', '3. Additional Context on the document and question asked'.\n"
                            + question,
                        }
                    ],
                },
            ]
            return self.model(messages).content
        except Exception as e:
            return str(e)